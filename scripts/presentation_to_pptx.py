#!/usr/bin/env python3
"""Genera presentazione_ia_produttiva.pptx dal contenuto della presentazione.

Per ogni slide della presentazione HTML/MD:
  1. estrae i contenuti,
  2. li mappa sul modello di slide corretto del template Qubit On Cloud,
  3. rigenera ogni schema come SVG (modulo svg_diagrams),
  4. lo inserisce nella slide come immagine SVG con fallback PNG.
"""

from __future__ import annotations

import copy
from io import BytesIO
from pathlib import Path

import cairosvg
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.opc.package import Part
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

import svg_diagrams as D

REPO = Path(__file__).resolve().parents[1]
TEMPLATE = REPO / "output" / "Qubit_On_Cloud_Template.pptx"
DST = REPO / "output" / "presentazione_ia_produttiva.pptx"

# Indici dei modelli-slide nel template
COVER, INDEX, DIVIDER, CONTENT, THREECOL = 0, 1, 2, 3, 4
VALUEPROP, TABLE, PROCESS, BIGSTAT, QUOTE, ROADMAP, CLOSING = 5, 6, 7, 8, 9, 10, 11

# Palette (senza '#') per python-pptx
C_BLUE_DARK = "003BA3"
C_BLUE = "4A8CFF"
C_INK = "1A1A2E"
C_MUTED = "4A4A5A"
C_WHITE = "FFFFFF"
C_SURFACE = "F5F7FA"
C_SURFACE2 = "E8EEF7"

# Namespace
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
ASVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"

_svg_counter = [0]


# --- Helper PPTX ------------------------------------------------------------
def clone_slide(prs, index):
    """Duplica un modello-slide del template come nuova slide."""
    src = prs.slides[index]
    blank = prs.slide_layouts[0]
    new = prs.slides.add_slide(blank)
    for ph in list(new.shapes):
        ph._element.getparent().remove(ph._element)
    for shp in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(shp._element))
    return new


def shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    raise KeyError(f"shape '{name}' non trovata")


def remove(slide, name):
    sh = shape(slide, name)
    sh._element.getparent().remove(sh._element)


def set_text(slide, name, text, size=None):
    """Sostituisce il testo preservando la formattazione del template.

    Supporta più righe (\\n) clonando il primo paragrafo.
    """
    tf = shape(slide, name).text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        p0.add_run()
    lines = str(text).split("\n")
    p0.runs[0].text = lines[0]
    for r in p0.runs[1:]:
        r._r.getparent().remove(r._r)
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    for line in lines[1:]:
        newp = copy.deepcopy(p0._p)
        tf._txBody.append(newp)
        from pptx.text.text import _Paragraph
        para = _Paragraph(newp, tf)
        para.runs[0].text = line
        for r in para.runs[1:]:
            r._r.getparent().remove(r._r)
    if size is not None:
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(size)


def fit_title(slide, name, text):
    """Imposta il titolo riducendo il corpo se troppo lungo."""
    n = len(text)
    size = 28 if n <= 40 else (24 if n <= 56 else 21)
    set_text(slide, name, text, size=size)


def add_svg(slide, prs, svg_str, left, top, width, height):
    """Inserisce un'immagine SVG (con fallback PNG) nella slide."""
    svg_bytes = svg_str.encode("utf-8")
    w_px = max(900, int(Emu(width).inches * 200))
    h_px = max(400, int(Emu(height).inches * 200))
    png = cairosvg.svg2png(bytestring=svg_bytes,
                           output_width=w_px, output_height=h_px)
    pic = slide.shapes.add_picture(BytesIO(png), left, top, width, height)

    _svg_counter[0] += 1
    part = slide.part
    partname = part.package.next_partname("/ppt/media/diagram%d.svg")
    svg_part = Part(partname, "image/svg+xml", part.package, svg_bytes)
    rId = part.relate_to(svg_part, RT.IMAGE)

    blip = pic._element.find(".//" + qn("a:blip"))
    extLst = etree.SubElement(blip, f"{{{A_NS}}}extLst")
    ext = etree.SubElement(extLst, f"{{{A_NS}}}ext")
    ext.set("uri", "{96DAC541-7B7A-43D3-8B79-37D633B846F1}")
    svgBlip = etree.SubElement(ext, f"{{{ASVG_NS}}}svgBlip")
    svgBlip.set(f"{{{R_NS}}}embed", rId)
    return pic


def add_diagram(slide, prs, svg_str):
    """Posiziona uno schema SVG nell'area contenuto del modello CONTENT."""
    add_svg(slide, prs, svg_str,
            Inches(0.45), Inches(1.18), Inches(9.1), Inches(4.05))


# --- Tabelle ----------------------------------------------------------------
def _style_cell(cell, text, size, color, bold, fill, align=PP_ALIGN.LEFT):
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor.from_string(fill)
    cell.margin_left = Inches(0.14)
    cell.margin_right = Inches(0.10)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = "Montserrat"
    r.font.color.rgb = RGBColor.from_string(color)


def table_slide(prs, title, headers, rows, widths, body_size=11):
    s = clone_slide(prs, TABLE)
    fit_title(s, "Text 0", title)
    remove(s, "Table 0")
    nrows, ncols = len(rows) + 1, len(headers)
    total_w = sum(widths)
    left, top = Inches(0.6), Inches(1.2)
    row_h = min(0.62, max(0.32, 3.9 / nrows))
    height = Inches(row_h * nrows)
    gt = s.shapes.add_table(nrows, ncols, left, top,
                            Inches(total_w), height).table
    gt.first_row = False
    gt.horz_banding = False
    for j, w in enumerate(widths):
        gt.columns[j].width = Inches(w)
    for r in range(nrows):
        gt.rows[r].height = Inches(row_h)
    for j, h in enumerate(headers):
        _style_cell(gt.cell(0, j), h, 12, C_WHITE, True, C_BLUE_DARK)
    for i, row in enumerate(rows, start=1):
        fill = C_WHITE if i % 2 else C_SURFACE
        for j, val in enumerate(row):
            bold = (j == 0)
            color = C_BLUE_DARK if j == 0 else C_INK
            _style_cell(gt.cell(i, j), val, body_size, color, bold, fill)
    return s


# --- Costruzione del deck ---------------------------------------------------
def build(prs):
    # 1 — TITOLO (cover)
    s = clone_slide(prs, COVER)
    set_text(s, "Text 1", "L'IA che Lavora")
    set_text(s, "Text 2",
             "Architettura, Dominio, Strategia\nCome trasformare un chatbot in un collega")
    set_text(s, "Text 3", "VALUATION ANALYST")
    set_text(s, "Text 4", "Caso reale · sistema multi-agente · 25 minuti")

    # 2 — IL PROBLEMA (content + piramide)
    s = clone_slide(prs, CONTENT)
    fit_title(s, "Text 0", "Il 90% usa l'IA come un motore di ricerca")
    remove(s, "Shape 3"); remove(s, "Text 4")
    add_diagram(s, prs, D.pyramid_adozione())

    # 3 — COLLABORATORE CREATIVO (tabella)
    table_slide(
        prs, "L'IA è un collaboratore creativo, non una calcolatrice",
        ["Calcolatrice (quello che già avete)",
         "Collaboratore creativo (quello che cambia tutto)"],
        [["Esegue una formula", "Sceglie quale formula usare e perché"],
         ["Produce un numero", "Chiede: questo numero ha senso?"],
         ["Fa quello che gli dici", "Propone quello che non avevi considerato"],
         ["Si ferma a un errore", "Spiega perché è un errore e suggerisce"],
         ["Scala in velocità", "Scala in giudizio"],
         ["Serve un programmatore", "Si adatta in linguaggio naturale"]],
        [4.4, 4.4], body_size=12)

    # 4 — I TRE PILASTRI (content + triangolo)
    s = clone_slide(prs, CONTENT)
    fit_title(s, "Text 0", "Tre pilastri: senza uno, il sistema non regge")
    remove(s, "Shape 3"); remove(s, "Text 4")
    add_diagram(s, prs, D.triangle_pilastri())

    # --- DIVIDER Pilastro 1
    s = clone_slide(prs, DIVIDER)
    set_text(s, "Text 2", "01")
    set_text(s, "Text 3", "Architettura")
    set_text(s, "Text 4", "Dal tool al sistema multi-agente")

    # 5 — UN TOOL NON SCALA (content + scala)
    s = clone_slide(prs, CONTENT)
    fit_title(s, "Text 0", "Un tool non scala: la scala della maturità")
    remove(s, "Shape 3"); remove(s, "Text 4")
    add_diagram(s, prs, D.staircase_maturita())

    # 6 — SISTEMA MULTI-AGENTE (content + hub-spoke + metriche)
    s = clone_slide(prs, CONTENT)
    fit_title(s, "Text 0", "Anatomia di un sistema multi-agente")
    remove(s, "Shape 3"); remove(s, "Text 4")
    add_diagram(s, prs, D.hub_spoke_sistema())

    # 7 — SKILL COME CONTRATTI (content + prima/dopo)
    s = clone_slide(prs, CONTENT)
    fit_title(s, "Text 0", "Le skill come contratti, non come prompt")
    remove(s, "Shape 3"); remove(s, "Text 4")
    add_diagram(s, prs, D.before_after_skill())

    # 8 — AUTONOMIA CONTROLLATA (content + gate timeline)
    s = clone_slide(prs, CONTENT)
    fit_title(s, "Text 0", "Autonomia controllata: sapere quando fermarsi")
    remove(s, "Shape 3"); remove(s, "Text 4")
    add_diagram(s, prs, D.gate_timeline())

    # 9 — SOTTO IL COFANO: STACK + FONTI DATI (chiude il Pilastro 1)
    table_slide(
        prs, "Sotto il cofano: stack tecnologico e fonti dati",
        ["Livello", "Tecnologia / Fonte", "Ruolo"],
        [["LLM", "Claude (Anthropic)", "Ragionamento, generazione, orchestrazione"],
         ["Runtime · prototipo", "Claude Code · MCP", "Esplorazione, tool use, sviluppo agenti"],
         ["Runtime · produzione", "Agent SDK · API", "Agenti in scope ristretto, modello via API"],
         ["Calcolo", "Python 3.11+", "Modelli finanziari · 348 test automatizzati"],
         ["Skill layer", "Markdown strutturato", "Workflow, vincoli, checkpoint"],
         ["Dati di mercato", "Massive · FactSet · Capital IQ · Refinitiv",
          "Prezzi, fondamentali, consensus"],
         ["Parametri settore", "Damodaran NYU Stern",
          "Beta, ERP, multipli, WACC settoriali"],
         ["Società italiane", "Bureau van Dijk (AIDA) · CONSOB · Borsa Italiana",
          "Bilanci, dati societari, depositi"],
         ["Versionamento", "Git + GitHub", "Audit trail, collaborazione"]],
        [1.9, 3.9, 3.0], body_size=10.5)

    # 10 — DAL BANCO DI LAVORO AL PRODOTTO (chiude il Pilastro 1)
    s = clone_slide(prs, CONTENT)
    fit_title(s, "Text 0", "Dal banco di lavoro al prodotto: prototipo → produzione")
    remove(s, "Shape 3"); remove(s, "Text 4")
    add_diagram(s, prs, D.prototype_to_production())

    # --- DIVIDER Pilastro 2
    s = clone_slide(prs, DIVIDER)
    set_text(s, "Text 2", "02")
    set_text(s, "Text 3", "Dominio")
    set_text(s, "Text 4", "La conoscenza che rende l'IA competente")

    # 9 — IL BRILLANTE INCOMPETENTE (content + due output)
    s = clone_slide(prs, CONTENT)
    fit_title(s, "Text 0", "Il brillante incompetente: senza dominio sbaglia")
    remove(s, "Shape 3"); remove(s, "Text 4")
    add_diagram(s, prs, D.two_outputs())

    # 10 — IL DOMINIO SI CODIFICA (content + pipeline conoscenza)
    s = clone_slide(prs, CONTENT)
    fit_title(s, "Text 0", "Il dominio si codifica: paper → codice → skill")
    remove(s, "Shape 3"); remove(s, "Text 4")
    add_diagram(s, prs, D.knowledge_pipeline())

    # 12 — COMPLIANCE COME DOMINIO (tabella) — il porting dei tool al framework UE
    table_slide(
        prs, "La compliance come dominio: il porting al framework europeo",
        ["Ambito", "Normativa", "Cosa fa l'IA"],
        [["Antitrust", "AGCM", "Verifica soglie fatturato (532M/32M €)"],
         ["Settori strategici", "Golden Power", "Flag settori D.L. 21/2012"],
         ["Mercati", "CONSOB", "Conformità Regolamento Emittenti"],
         ["Adeguatezza", "MiFID II", "Profiling, target market, ESG"],
         ["Antiriciclaggio", "D.Lgs. 231/2007", "3 livelli di adeguata verifica"],
         ["Sostenibilità", "SFDR", "Classificazione art. 6/8/9"],
         ["Resilienza", "DORA", "Check Reg. UE 2022/2554"],
         ["Fiscalità", "IRES+IRAP", "Parametri automatici per l'Italia"]],
        [2.0, 2.4, 4.4], body_size=11)

    # 13 — ARRICCHIMENTO DEL DOMINIO (content + piramide) — la storia del progetto
    s = clone_slide(prs, CONTENT)
    fit_title(s, "Text 0", "Arricchire il dominio: dallo strumento alla suite integrata")
    remove(s, "Shape 3"); remove(s, "Text 4")
    add_diagram(s, prs, D.domain_enrichment())

    # --- DIVIDER Pilastro 3
    s = clone_slide(prs, DIVIDER)
    set_text(s, "Text 2", "03")
    set_text(s, "Text 3", "Strategia")
    set_text(s, "Text 4", "Il percorso di adozione progressiva")

    # 13 — NON SERVE UN BIG BANG (content + tre binari)
    s = clone_slide(prs, CONTENT)
    fit_title(s, "Text 0", "Non serve un big bang: adozione progressiva")
    remove(s, "Shape 3"); remove(s, "Text 4")
    add_diagram(s, prs, D.three_rails())

    # 14 — L'UOMO NEL LOOP (content + matrice)
    s = clone_slide(prs, CONTENT)
    fit_title(s, "Text 0", "L'uomo nel loop: dove metterlo e dove toglierlo")
    remove(s, "Shape 3"); remove(s, "Text 4")
    add_diagram(s, prs, D.matrix_2x2())

    # 15 — MISURARE IL VALORE (tabella)
    table_slide(
        prs, "Misurare il valore: metriche che contano",
        ["Metriche vanity (inutili)", "Metriche di valore (utili)"],
        [["Numero di prompt al giorno", "Tempo da richiesta a output utilizzabile"],
         ["Token consumati", "Tasso di rework (output accettati al 1º giro)"],
         ["Utenti attivi", "Copertura dei controlli (check eseguiti)"],
         ["Soddisfazione generica", "Riproducibilità (stesso input = stesso output)"],
         ["Costo per query", "Audit trail (tracciabilità delle decisioni)"]],
        [4.4, 4.4], body_size=12)

    # --- DIVIDER Chiusura
    s = clone_slide(prs, DIVIDER)
    set_text(s, "Text 2", "★")
    set_text(s, "Text 3", "Chiusura")
    set_text(s, "Text 4", "I tre cardini e le azioni concrete")

    # 16 — I TRE CARDINI / RECAP (tre colonne)
    s = clone_slide(prs, THREECOL)
    fit_title(s, "Text 0", "I tre cardini: senza uno, il sistema non regge")
    set_text(s, "Text 5", "ARCHITETTURA")
    set_text(s, "Text 6",
             "Non un prompt, un sistema. Agenti specializzati, skill con "
             "garanzie, autonomia controllata.\n\nSenza: IA brillante ma caotica.")
    set_text(s, "Text 9", "DOMINIO")
    set_text(s, "Text 10",
             "Non dati, giudizio. Paper, formule, vincoli ed errori comuni "
             "codificati.\n\nSenza: IA veloce ma pericolosa.")
    set_text(s, "Text 13", "STRATEGIA")
    set_text(s, "Text 14",
             "Non big bang, percorso. Partire da dove il rischio è basso e il "
             "valore visibile.\n\nSenza: pilota che non finisce mai.")

    # 17 — CALL TO ACTION (value proposition)
    s = clone_slide(prs, VALUEPROP)
    fit_title(s, "Text 0", "Call to action: 3 mosse per lunedì mattina")
    set_text(s, "Text 4", "Costruite\nsistemi,\nnon prompt")
    set_text(s, "Text 5",
             "Smettete di chiedere all'IA di fare cose. Iniziate a costruire "
             "sistemi che sanno farle.")
    set_text(s, "Text 6", "1 · Scrivete un workflow come skill")
    set_text(s, "Text 7",
             "Vincoli, step, checkpoint, errori comuni. Non serve codice: basta un documento.")
    set_text(s, "Text 8", "2 · Mappate i decision gate")
    set_text(s, "Text 9",
             "Dove l'IA procede da sola e dove si ferma. Usate la matrice rischio/reversibilità.")
    set_text(s, "Text 10", "3 · Partite piccoli e misurate")
    set_text(s, "Text 11",
             "Un agente, una skill, un workflow. Misurate gli output utilizzabili, non i prompt.")

    # 18 — GRAZIE (closing)
    s = clone_slide(prs, CLOSING)
    set_text(s, "Text 2", "Grazie")
    set_text(s, "Text 3", "Domande?")
    set_text(s, "Text 4",
             "Il futuro non è l'IA che sostituisce le persone.\n"
             "È l'IA che le rende più brave nel loro mestiere.")
    set_text(s, "Text 5", "VALUATION ANALYST")

    # --- DIVIDER Backup
    s = clone_slide(prs, DIVIDER)
    set_text(s, "Text 2", "+")
    set_text(s, "Text 3", "Backup")
    set_text(s, "Text 4", "Approfondimenti per le domande")

    # B1 — PIPELINE CREDIT RISK (processo 5 fasi)
    s = clone_slide(prs, PROCESS)
    fit_title(s, "Text 0", "Backup · Dal bilancio al rating in 5 step")
    set_text(s, "Text 4", "Dataset")
    set_text(s, "Text 6", "48 aziende, invarianti di bilancio verificati")
    set_text(s, "Text 8", "BMS")
    set_text(s, "Text 10", "Media di 35 peer: EBITDA 14%, leva D/TA 42%")
    set_text(s, "Text 12", "Differenziale")
    set_text(s, "Text 14", "Target vs settore: margine +3pp, 6/8 favorevoli")
    set_text(s, "Text 16", "DCF coerente")
    set_text(s, "Text 18", "2 stadi, 7 check PASS, EV 38M, TV peso 64%")
    set_text(s, "Text 20", "Monte Carlo")
    set_text(s, "Text 21", "20k sim → PD 2,8%, rating BB+, EL 1,2M")

    # B3 — ROI DELL'APPROCCIO (tabella)
    table_slide(
        prs, "Backup · ROI dell'approccio",
        ["Investimento — una tantum (~1 settimana)", "Ritorno — ricorrente"],
        [["Codifica dominio: 2-3 giorni per verticale",
          "Tempo analisi −80% (2-3 giorni → 2-3 ore)"],
         ["Architettura agenti: 1-2 giorni di setup",
          "7+ check sempre eseguiti, zero dimenticanze"],
         ["Test e validazione: 1 giorno",
          "Riproducibilità garantita, audit-friendly"],
         ["Totale setup ~1 settimana",
          "Onboarding junior: da mesi a giorni (skill = formazione)"]],
        [4.4, 4.4], body_size=12)


def drop_template_slides(prs, n):
    """Rimuove i primi n modelli-slide del template, lasciando solo il deck."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst)[:n]:
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def main():
    prs = Presentation(str(TEMPLATE))
    n_template = len(prs.slides)
    build(prs)
    drop_template_slides(prs, n_template)
    prs.save(str(DST))
    print(f"Presentazione PPTX generata: {DST}")
    print(f"Slide totali: {len(prs.slides)}  ·  diagrammi SVG: {_svg_counter[0]}")
    print(f"Dimensione: {DST.stat().st_size / 1024:.0f} KB")


if __name__ == "__main__":
    main()
